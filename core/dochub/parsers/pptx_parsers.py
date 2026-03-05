from abc import ABC

import pptx
from pptx.parts.image import Image
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture
from pptx.opc.constants import CONTENT_TYPE

from loguru import logger

from dochub.parsers.base import BaseDocumentParser
from dochub.parsers.image_parsers import OCRModelEnum
from dochub.schemas import Document, Chunk, DataType, Param, ChunkType
from dochub.utils import ocr_utils
from utils.i18n import I18NString, Language


class BasePPTXParser(BaseDocumentParser, ABC):
    target_content_type = [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/wps-office.pptx"
    ]
    target_file_ext = "pptx"
    pptx_dir_prefix = "/"

    def _validate_document(self, target: Document):
        assert target.content_type in self.target_content_type, "文档类型不匹配！期望的文档类型：{}，得到的文档类型：{}".format(
            str(self.target_content_type), target.content_type)


class GeneralPPTXParser(BasePPTXParser):
    name = I18NString({
        Language.ZH: "通用",
        Language.EN: "General",
    })

    params = [
        Param(name="ocr_enable", display_name="是否采用OCR", required=True, data_type=DataType.BOOLEAN, default=True),
        Param(name="ocr_model", display_name="OCR模型", required=True, data_type=DataType.STRING,
              default=OCRModelEnum.RapidOCR),
    ]

    def __init__(self, target: Document, **kwargs):
        super().__init__(target, **kwargs)
        self.ocr_enable = kwargs.get("ocr_enable", False)
        self.ocr_model = kwargs.get("ocr_model", OCRModelEnum.RapidOCR)

    def _parse_impl(self) -> None:
        prs = pptx.Presentation(self.target.physical_path)

        total_cnt = len(prs.slides)
        for i, slide in enumerate(self._progress_wrapper(prs.slides, total_cnt)):
            for shape in slide.shapes:
                if isinstance(shape, GraphicFrame):
                    # 可能是 table, chart, smart art, and media objects.
                    # 先只处理 table
                    self._parse_table(shape, i + 1)
                else:
                    # 可能是 AutoShape 或 GroupShape
                    self._parse_text(shape, i + 1)
        self._report_progress(100)

    def _parse_text(self, shape, slide_num: int) -> None:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = "".join([run.text for run in paragraph.runs])
                if len(text.strip()) == 0:
                    continue
                chunk = Chunk(content=text, type=ChunkType.TEXT, metadata={"pages": str(slide_num)})
                self._append_content_chunks(chunk)
            return
        if isinstance(shape, Picture) and self.ocr_enable:
            self._parse_picture_text(shape, slide_num)
            return
        if isinstance(shape, GroupShape):
            for sp in shape.shapes:
                self._parse_text(sp, slide_num)

    def _parse_picture_text(self, shape: Picture, slide_num: int) -> None:
        """
        使用OCR解析出图片中的文字
        """
        try:
            image: Image = shape.image
        except ValueError as e:
            logger.warning("This picture does not have an image", e)
            return
        if image.content_type == CONTENT_TYPE.GIF:
            # TODO: 如何解析 GIF
            return
        if shape._pic.spPr.xpath(".//a:scene3d"):
            # 被水平翻折的图片暂不处理
            return
        pic_content = image.blob
        if (image.content_type == CONTENT_TYPE.X_WMF or
                image.content_type == CONTENT_TYPE.X_EMF):
            # TODO: 将emf转为png来解析
            return
        try:
            self._append_content_chunks(*ocr_utils.ocr_parse_grouped(pic_content, slide_num, self.ocr_model))
        except Exception as e:
            logger.exception(f"ocr failed in image {image.filename}", e)

    def _parse_table(self, gf: GraphicFrame, slide_num: int) -> None:
        """
        只解析 ``gf`` 中的表格，如果有合并单元格，只展示最左上的单元格的内容，被合并的单元格的内容是""
        """
        if gf.has_table:
            table = gf.table
            m, n = len(table.rows), len(table.columns)
            text = ""
            for i in range(0, m):
                text += "|"
                for j in range(0, n):
                    text += table.cell(i, j).text
                    text += "|"
                text += "\n"
            tbl_metadata = {
                "pages": str(slide_num),
                "rows": len(table.rows),
                "columns": len(table.columns),
            }
            chunk = Chunk(content=text, type=ChunkType.TABLE, metadata=tbl_metadata)
            self._append_content_chunks(chunk)
